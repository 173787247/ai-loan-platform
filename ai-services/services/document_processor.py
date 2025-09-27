"""
文档处理服务

@author AI Loan Platform Team
@version 1.0.0
"""

import os
import json
import logging
from typing import Dict, Any, Optional, List
import pytesseract
from PIL import Image
import PyPDF2
from docx import Document
import cv2
import numpy as np
import pandas as pd
from pptx import Presentation
import openpyxl
import xlrd
import xlwt
import markdown
from bs4 import BeautifulSoup
import magic
from loguru import logger
from .advanced_ocr import advanced_ocr_service, OCREngine

class DocumentProcessor:
    """文档处理服务类"""
    
    def __init__(self):
        self.supported_types = [
            # 文档格式
            'pdf', 'doc', 'docx', 'rtf', 'txt', 'md',
            # 表格格式
            'xls', 'xlsx', 'csv',
            # 演示文稿格式
            'ppt', 'pptx',
            # 图片格式
            'jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif',
            # 网页格式
            'html', 'htm'
        ]
        self.logger = logger
        
    async def process_document(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """
        处理文档（异步版本）
        
        Args:
            file_path: 文件路径
            file_type: 文件类型
            
        Returns:
            处理结果
        """
        try:
            self.logger.info(f"开始处理文档: {file_path}")
            
            # 检查文件类型
            if not self._is_supported_type(file_type):
                raise ValueError(f"不支持的文件类型: {file_type}")
            
            # 提取文本（异步）
            text = await self._extract_text_async(file_path, file_type)
            
            # 分类文档
            doc_type = self._classify_document(text)
            
            # 提取结构化数据
            structured_data = self._extract_structured_data(text, doc_type)
            
            # 验证数据完整性
            validation_result = self._validate_data(structured_data, doc_type)
            
            result = {
                "text": text,
                "document_type": doc_type,
                "structured_data": structured_data,
                "validation": validation_result,
                "file_info": {
                    "path": file_path,
                    "type": file_type,
                    "size": os.path.getsize(file_path)
                }
            }
            
            self.logger.info(f"文档处理完成: {file_path}")
            return result
            
        except Exception as e:
            self.logger.error(f"文档处理失败: {file_path}, 错误: {str(e)}")
            raise
    
    def _is_supported_type(self, file_type: str) -> bool:
        """检查文件类型是否支持"""
        return any(file_type.lower().endswith(ext) for ext in self.supported_types)
    
    async def _extract_text_async(self, file_path: str, file_type: str) -> str:
        """提取文档文本（异步版本）"""
        try:
            file_ext = file_type.lower()
            
            # PDF文档
            if file_ext.endswith('pdf'):
                return self._extract_pdf_text(file_path)
            
            # Word文档
            elif file_ext.endswith(('doc', 'docx')):
                return self._extract_word_text(file_path)
            
            # Excel表格
            elif file_ext.endswith(('xls', 'xlsx')):
                return self._extract_excel_text(file_path)
            
            # PowerPoint演示文稿
            elif file_ext.endswith(('ppt', 'pptx')):
                return self._extract_ppt_text(file_path)
            
            # 图片文件（使用高级OCR）
            elif file_ext.endswith(('jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif')):
                return await self._extract_image_text(file_path)
            
            # 其他文件类型
            else:
                return self._extract_text(file_path, file_type)
                
        except Exception as e:
            self.logger.error(f"异步文本提取失败: {file_path}, 错误: {str(e)}")
            return ""
    
    def _extract_text(self, file_path: str, file_type: str) -> str:
        """提取文档文本（同步版本）"""
        try:
            file_ext = file_type.lower()
            
            # PDF文档
            if file_ext.endswith('pdf'):
                return self._extract_pdf_text(file_path)
            
            # Word文档
            elif file_ext.endswith(('doc', 'docx')):
                return self._extract_word_text(file_path)
            
            # Excel表格
            elif file_ext.endswith(('xls', 'xlsx')):
                return self._extract_excel_text(file_path)
            
            # PowerPoint演示文稿
            elif file_ext.endswith(('ppt', 'pptx')):
                return self._extract_ppt_text(file_path)
            
            # CSV文件
            elif file_ext.endswith('csv'):
                return self._extract_csv_text(file_path)
            
            # 文本文件
            elif file_ext.endswith(('txt', 'md', 'rtf')):
                return self._extract_text_file(file_path)
            
            # HTML文件
            elif file_ext.endswith(('html', 'htm')):
                return self._extract_html_text(file_path)
            
            # 图片文件（OCR）
            elif file_ext.endswith(('jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif')):
                return self._extract_image_text(file_path)
            
            else:
                raise ValueError(f"不支持的文件类型: {file_type}")
        except Exception as e:
            self.logger.error(f"文本提取失败: {file_path}, 错误: {str(e)}")
            raise
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """提取PDF文本和图片OCR - 使用多种方法确保最佳效果"""
        text = ""
        
        # 方法1: 尝试pdfplumber (通常效果最好)
        try:
            import pdfplumber
            # 设置中文字体支持
            import os
            os.environ['FONTCONFIG_PATH'] = '/etc/fonts'
            
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                    
                    # 提取图片并进行OCR
                    images = page.images
                    if images:
                        self.logger.info(f"发现 {len(images)} 个图片，开始OCR处理")
                        for i, img in enumerate(images):
                            try:
                                # 获取图片边界框
                                bbox = [img.get('x0', 0), img.get('y0', 0), img.get('x1', 0), img.get('y1', 0)]
                                if bbox[2] > bbox[0] and bbox[3] > bbox[1]:  # 确保边界框有效
                                    # 提取图片区域
                                    img_obj = page.within_bbox(bbox).to_image()
                                    if img_obj:
                                        # 转换为PIL图片
                                        pil_img = img_obj.original
                                        # OCR识别
                                        img_text = pytesseract.image_to_string(pil_img, lang='chi_sim+eng')
                                        if img_text.strip():
                                            text += f"\n[图片{i+1}OCR内容]:\n{img_text.strip()}\n"
                                            self.logger.info(f"图片{i+1} OCR成功: {len(img_text)}字符")
                                        else:
                                            self.logger.info(f"图片{i+1} OCR结果为空")
                            except Exception as e:
                                self.logger.warning(f"图片{i+1} OCR失败: {e}")
            
            if text.strip():
                self.logger.info(f"pdfplumber成功提取PDF文本和图片OCR: {len(text)}字符")
                return text.strip()
        except Exception as e:
            self.logger.warning(f"pdfplumber提取失败: {e}")
        
        # 方法2: 尝试PyMuPDF (fitz) - 支持图片提取
        try:
            import fitz
            # 设置中文字体支持
            import os
            os.environ['FONTCONFIG_PATH'] = '/etc/fonts'
            
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_text = page.get_text()
                if page_text:
                    text += page_text + "\n"
                
                # 提取图片并进行OCR
                image_list = page.get_images()
                if image_list:
                    self.logger.info(f"页面{page_num+1}发现 {len(image_list)} 个图片，开始OCR处理")
                    for img_index, img in enumerate(image_list):
                        try:
                            # 获取图片
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            if pix.n - pix.alpha < 4:  # 确保不是CMYK
                                # 转换为PIL图片
                                img_data = pix.tobytes("png")
                                from io import BytesIO
                                pil_img = Image.open(BytesIO(img_data))
                                
                                # OCR识别
                                img_text = pytesseract.image_to_string(pil_img, lang='chi_sim+eng')
                                if img_text.strip():
                                    text += f"\n[页面{page_num+1}图片{img_index+1}OCR内容]:\n{img_text.strip()}\n"
                                    self.logger.info(f"页面{page_num+1}图片{img_index+1} OCR成功: {len(img_text)}字符")
                            pix = None
                        except Exception as e:
                            self.logger.warning(f"页面{page_num+1}图片{img_index+1} OCR失败: {e}")
            
            doc.close()
            if text.strip():
                self.logger.info(f"PyMuPDF成功提取PDF文本和图片OCR: {len(text)}字符")
                return text.strip()
        except Exception as e:
            self.logger.warning(f"PyMuPDF提取失败: {e}")
        
        # 方法3: 回退到PyPDF2
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            if text.strip():
                self.logger.info(f"PyPDF2成功提取PDF文本: {len(text)}字符")
                return text.strip()
        except Exception as e:
            self.logger.warning(f"PyPDF2提取失败: {e}")
        
        self.logger.error(f"所有PDF提取方法都失败了: {file_path}")
        return ""
    
    def _extract_word_text(self, file_path: str) -> str:
        """提取Word文档文本"""
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    
    def _extract_image_text(self, file_path: str) -> str:
        """提取图片文本（OCR）"""
        # 预处理图片
        image = cv2.imread(file_path)
        if image is None:
            raise ValueError(f"无法读取图片: {file_path}")
        
        # 图片预处理
        processed_image = self._preprocess_image(image)
        
        # OCR识别
        text = pytesseract.image_to_string(processed_image, lang='chi_sim+eng')
        return text.strip()
    
    def _preprocess_image(self, image: np.ndarray) -> np.ndarray:
        """预处理图片以提高OCR准确率"""
        # 转换为灰度图
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        
        # 去噪
        denoised = cv2.medianBlur(gray, 3)
        
        # 二值化
        _, binary = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # 形态学操作
        kernel = np.ones((1, 1), np.uint8)
        processed = cv2.morphologyEx(binary, cv2.MORPH_CLOSE, kernel)
        
        return processed
    
    def _classify_document(self, text: str) -> str:
        """分类文档类型"""
        text_lower = text.lower()
        
        # 营业执照关键词
        business_license_keywords = ['营业执照', '统一社会信用代码', '法定代表人', '注册资本']
        if any(keyword in text for keyword in business_license_keywords):
            return 'business_license'
        
        # 财务报表关键词
        financial_keywords = ['资产负债表', '利润表', '现金流量表', '营业收入', '净利润']
        if any(keyword in text for keyword in financial_keywords):
            return 'financial_statement'
        
        # 银行流水关键词
        bank_statement_keywords = ['银行流水', '交易明细', '账户余额', '转账记录']
        if any(keyword in text for keyword in bank_statement_keywords):
            return 'bank_statement'
        
        # 税务报表关键词
        tax_keywords = ['纳税申报表', '增值税', '企业所得税', '税务登记']
        if any(keyword in text for keyword in tax_keywords):
            return 'tax_statement'
        
        # 其他文档
        return 'other'
    
    def _extract_structured_data(self, text: str, doc_type: str) -> Dict[str, Any]:
        """提取结构化数据"""
        if doc_type == 'business_license':
            return self._extract_business_license_data(text)
        elif doc_type == 'financial_statement':
            return self._extract_financial_data(text)
        elif doc_type == 'bank_statement':
            return self._extract_bank_statement_data(text)
        else:
            return {"raw_text": text}
    
    def _extract_business_license_data(self, text: str) -> Dict[str, Any]:
        """提取营业执照数据"""
        data = {}
        
        # 提取公司名称
        import re
        company_name_pattern = r'名称[：:]\s*([^\n\r]+)'
        match = re.search(company_name_pattern, text)
        if match:
            data['company_name'] = match.group(1).strip()
        
        # 提取统一社会信用代码
        credit_code_pattern = r'统一社会信用代码[：:]\s*([A-Z0-9]{18})'
        match = re.search(credit_code_pattern, text)
        if match:
            data['credit_code'] = match.group(1)
        
        # 提取法定代表人
        legal_person_pattern = r'法定代表人[：:]\s*([^\n\r]+)'
        match = re.search(legal_person_pattern, text)
        if match:
            data['legal_person'] = match.group(1).strip()
        
        # 提取注册资本
        capital_pattern = r'注册资本[：:]\s*([^\n\r]+)'
        match = re.search(capital_pattern, text)
        if match:
            data['registered_capital'] = match.group(1).strip()
        
        return data
    
    def _extract_financial_data(self, text: str) -> Dict[str, Any]:
        """提取财务报表数据"""
        data = {}
        
        # 提取营业收入
        revenue_pattern = r'营业收入[：:]\s*([0-9,]+\.?[0-9]*)'
        match = re.search(revenue_pattern, text)
        if match:
            data['revenue'] = float(match.group(1).replace(',', ''))
        
        # 提取净利润
        profit_pattern = r'净利润[：:]\s*([0-9,]+\.?[0-9]*)'
        match = re.search(profit_pattern, text)
        if match:
            data['net_profit'] = float(match.group(1).replace(',', ''))
        
        # 提取总资产
        assets_pattern = r'总资产[：:]\s*([0-9,]+\.?[0-9]*)'
        match = re.search(assets_pattern, text)
        if match:
            data['total_assets'] = float(match.group(1).replace(',', ''))
        
        return data
    
    def _extract_bank_statement_data(self, text: str) -> Dict[str, Any]:
        """提取银行流水数据"""
        data = {}
        
        # 提取账户余额
        balance_pattern = r'余额[：:]\s*([0-9,]+\.?[0-9]*)'
        match = re.search(balance_pattern, text)
        if match:
            data['balance'] = float(match.group(1).replace(',', ''))
        
        # 提取交易笔数
        transaction_count = len(re.findall(r'\d{4}-\d{2}-\d{2}', text))
        data['transaction_count'] = transaction_count
        
        return data
    
    def _validate_data(self, data: Dict[str, Any], doc_type: str) -> Dict[str, Any]:
        """验证数据完整性"""
        validation = {
            "is_valid": True,
            "missing_fields": [],
            "warnings": []
        }
        
        if doc_type == 'business_license':
            required_fields = ['company_name', 'credit_code', 'legal_person']
            for field in required_fields:
                if field not in data or not data[field]:
                    validation["missing_fields"].append(field)
                    validation["is_valid"] = False
        
        elif doc_type == 'financial_statement':
            required_fields = ['revenue', 'net_profit']
            for field in required_fields:
                if field not in data or data[field] is None:
                    validation["missing_fields"].append(field)
                    validation["is_valid"] = False
        
        return validation
    
    def get_status(self) -> Dict[str, Any]:
        """获取服务状态"""
        return {
            "status": "running",
            "supported_types": self.supported_types,
            "version": "1.0.0"
        }
    
    def _extract_excel_text(self, file_path: str) -> str:
        """提取Excel表格文本"""
        try:
            # 使用pandas读取Excel文件
            if file_path.endswith('.xlsx'):
                df = pd.read_excel(file_path, sheet_name=None)  # 读取所有工作表
            else:  # .xls
                df = pd.read_excel(file_path, sheet_name=None, engine='xlrd')
            
            text_parts = []
            for sheet_name, sheet_df in df.items():
                text_parts.append(f"工作表: {sheet_name}")
                # 将DataFrame转换为文本
                text_parts.append(sheet_df.to_string(index=False))
                text_parts.append("\n")
            
            return "\n".join(text_parts)
        except Exception as e:
            self.logger.error(f"Excel文件读取失败: {e}")
            return ""
    
    def _extract_ppt_text(self, file_path: str) -> str:
        """提取PowerPoint演示文稿文本"""
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"幻灯片 {i}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
                text_parts.append("\n")
            
            return "\n".join(text_parts)
        except Exception as e:
            self.logger.error(f"PowerPoint文件读取失败: {e}")
            return ""
    
    def _extract_csv_text(self, file_path: str) -> str:
        """提取CSV文件文本"""
        try:
            df = pd.read_csv(file_path, encoding='utf-8')
            return df.to_string(index=False)
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                df = pd.read_csv(file_path, encoding='gbk')
                return df.to_string(index=False)
            except Exception as e:
                self.logger.error(f"CSV文件读取失败: {e}")
                return ""
        except Exception as e:
            self.logger.error(f"CSV文件读取失败: {e}")
            return ""
    
    def _extract_text_file(self, file_path: str) -> str:
        """提取文本文件内容"""
        try:
            # 尝试UTF-8编码
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # 尝试GBK编码
            try:
                with open(file_path, 'r', encoding='gbk') as file:
                    return file.read()
            except UnicodeDecodeError:
                # 尝试latin-1编码
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
        except Exception as e:
            self.logger.error(f"文本文件读取失败: {e}")
            return ""
    
    def _extract_html_text(self, file_path: str) -> str:
        """提取HTML文件文本"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            # 使用BeautifulSoup解析HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 移除脚本和样式标签
            for script in soup(["script", "style"]):
                script.decompose()
            
            # 获取文本内容
            text = soup.get_text()
            
            # 清理文本
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text
        except Exception as e:
            self.logger.error(f"HTML文件读取失败: {e}")
            return ""
    
    async def _extract_image_text(self, file_path: str) -> str:
        """提取图片文本（使用高级OCR）"""
        try:
            # 使用高级OCR服务
            ocr_results = await advanced_ocr_service.recognize_text(
                file_path, 
                engines=[OCREngine.PADDLEOCR, OCREngine.TESSERACT],
                language="chi_sim+eng"
            )
            
            if not ocr_results:
                self.logger.warning("所有OCR引擎都失败了，尝试传统方法")
                return self._fallback_ocr(file_path)
            
            # 获取最佳结果
            best_result = await advanced_ocr_service.get_best_result(ocr_results)
            
            if best_result and best_result.text.strip():
                self.logger.info(f"OCR识别成功: {best_result.engine}, 置信度: {best_result.confidence:.3f}")
                return best_result.text.strip()
            else:
                self.logger.warning("高级OCR未识别到文本，尝试传统方法")
                return self._fallback_ocr(file_path)
            
        except Exception as e:
            self.logger.error(f"高级OCR处理失败: {file_path}, 错误: {str(e)}")
            return self._fallback_ocr(file_path)
    
    def _fallback_ocr(self, file_path: str) -> str:
        """传统OCR方法作为备用"""
        try:
            # 读取图片
            image = cv2.imread(file_path)
            if image is None:
                raise ValueError(f"无法读取图片: {file_path}")
            
            # 图片预处理
            processed_images = self._preprocess_image_for_ocr(image)
            
            # 对每个处理后的图片进行OCR
            all_text = []
            for processed_img in processed_images:
                try:
                    # 使用多种OCR配置
                    text = self._perform_ocr(processed_img)
                    if text.strip():
                        all_text.append(text.strip())
                except Exception as e:
                    self.logger.warning(f"传统OCR处理失败: {e}")
                    continue
            
            # 合并所有文本
            final_text = "\n".join(all_text)
            
            # 后处理文本
            final_text = self._postprocess_ocr_text(final_text)
            
            return final_text
            
        except Exception as e:
            self.logger.error(f"传统OCR处理失败: {file_path}, 错误: {str(e)}")
            return ""
    
    def _preprocess_image_for_ocr(self, image: np.ndarray) -> List[np.ndarray]:
        """为OCR预处理图片，返回多个处理版本"""
        processed_images = []
        
        # 原始图片
        processed_images.append(image)
        
        # 转换为灰度图
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        processed_images.append(gray)
        
        # 高斯模糊去噪
        denoised = cv2.GaussianBlur(gray, (5, 5), 0)
        processed_images.append(denoised)
        
        # 自适应阈值二值化
        binary = cv2.adaptiveThreshold(
            gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
        )
        processed_images.append(binary)
        
        # 形态学操作
        kernel = np.ones((2, 2), np.uint8)
        morph = cv2.morphologyEx(binary, cv2.MORPH_CLOSE, kernel)
        processed_images.append(morph)
        
        # 边缘检测
        edges = cv2.Canny(gray, 50, 150)
        processed_images.append(edges)
        
        return processed_images
    
    def _perform_ocr(self, image: np.ndarray) -> str:
        """执行OCR识别"""
        try:
            # 配置OCR参数
            config = '--oem 3 --psm 6 -l chi_sim+eng'
            
            # 执行OCR
            text = pytesseract.image_to_string(image, config=config)
            
            return text
        except Exception as e:
            self.logger.warning(f"OCR识别失败: {e}")
            return ""
    
    def _postprocess_ocr_text(self, text: str) -> str:
        """后处理OCR文本"""
        if not text:
            return ""
        
        # 移除多余的空白字符
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            if line and len(line) > 1:  # 过滤掉单字符和空行
                cleaned_lines.append(line)
        
        # 合并相邻的短行
        merged_lines = []
        current_line = ""
        
        for line in cleaned_lines:
            if len(line) < 10 and current_line:
                # 短行，合并到当前行
                current_line += " " + line
            else:
                if current_line:
                    merged_lines.append(current_line)
                current_line = line
        
        if current_line:
            merged_lines.append(current_line)
        
        return "\n".join(merged_lines)
    
    def batch_process_documents(self, file_paths: List[str]) -> Dict[str, Any]:
        """批量处理文档"""
        results = {}
        success_count = 0
        error_count = 0
        
        for file_path in file_paths:
            try:
                file_type = os.path.splitext(file_path)[1][1:]  # 获取文件扩展名
                result = self.process_document(file_path, file_type)
                results[file_path] = result
                success_count += 1
            except Exception as e:
                self.logger.error(f"批量处理失败: {file_path}, 错误: {str(e)}")
                results[file_path] = {"error": str(e)}
                error_count += 1
        
        return {
            "results": results,
            "summary": {
                "total": len(file_paths),
                "success": success_count,
                "error": error_count
            }
        }
